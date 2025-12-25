#!/usr/bin/env python3
"""
EAC-DTW Condensed Presentation Generator (8-9 Slides)
Creates a focused presentation aligned with poster content
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_condensed_presentation():
    """Create an 8-slide condensed presentation aligned with poster content"""

    # Create presentation object
    prs = Presentation()

    # Slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content
    section_header_layout = prs.slide_layouts[2]  # Section header

    # Color scheme (Pace University colors)
    PACE_BLUE = RGBColor(0, 51, 102)
    PACE_GREEN = RGBColor(0, 128, 0)
    PACE_RED = RGBColor(178, 34, 34)

    # Image paths
    image_dir = os.getcwd()
    img_paths = {
        'dtw_comparison_3': os.path.join(image_dir, 'dtw_comparison_3.png'),
        'cost_matrix_comparison': os.path.join(image_dir, 'cost_matrix_comparison.png'),
        'EG_DTW_Performance_2025': os.path.join(image_dir, 'EG_DTW_Performance_2025.png'),
        'singularities_comparision': os.path.join(image_dir, 'singularities_comparision.png'),
        'ECG_Simulation_with_noise': os.path.join(image_dir, 'ECG_Simulation_with_noise.png')
    }

    # ===== SLIDE 1: TITLE SLIDE =====
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    title.text = "EAC-DTW: Entropy-Adaptive Constraint Dynamic Time Warping Framework for Quantifiably Trustworthy ECG Classification"

    subtitle = slide.placeholders[1]
    subtitle.text = """Fnu Ashutosh (an05893n@pace.edu)
Shivam Jha (sj34101n@pace.edu)
Faculty Advisor: Dr. Sung-Hyuk Cha (scha@pace.edu)

Seidenberg School of Computer Science and Information Systems
Pace University, New York, NY

December 5, 2025"""

    # ===== SLIDE 2: CLINICAL CONTEXT & PROBLEM =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'The ECG Classification Challenge'
    tf = body_shape.text_frame
    tf.text = '''Cardiovascular diseases (CVDs) remain the leading cause of mortality globally, necessitating automated ECG analysis for arrhythmia detection.

The DTW Dilemma:
• Standard DTW enables elastic alignment but creates "pathological warping"
• Fixed constraints (Sakoe-Chiba band) are too rigid for complex regions, too permissive for noise
• ECG signals have heterogeneous complexity: high in QRS complexes, low in isoelectric segments

Our Solution: Entropy-Adaptive Constraint DTW (EAC-DTW)
• Uses local Shannon entropy to dynamically adjust warping constraints
• Tight constraints in noisy/flat regions, flexible in complex morphological areas
• Maintains alignment accuracy while preventing spurious matches'''

    # Add image if available
    if os.path.exists(img_paths['ECG_Simulation_with_noise']):
        left = Inches(4.5)
        top = Inches(2.0)
        width = Inches(3.5)
        height = Inches(2.5)
        slide.shapes.add_picture(img_paths['ECG_Simulation_with_noise'], left, top, width, height)

    # ===== SLIDE 3: RELATED WORK =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Related Work: Fixed Constraint Limitations'
    tf = body_shape.text_frame
    tf.text = '''Sakoe-Chiba Band (1978):
• Restricts warping to fixed diagonal band |i-j| ≤ R
• Reduces complexity from O(N²) to O(NR)
• Industry standard but "one-size-fits-all"

Critical Limitations:
❌ Cannot adapt to ECG signal heterogeneity
❌ Too rigid for QRS complexes (needs flexibility)
❌ Too permissive in isoelectric regions (allows noise alignment)

Other Approaches:
• Derivative DTW: Amplifies noise
• Soft-DTW: Computationally expensive
• Fixed constraints: Ignore local signal complexity

Key Insight: ECG signals need position-dependent constraints based on local information content.'''

    # ===== SLIDE 4: EAC-DTW METHODOLOGY =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'EAC-DTW: Adaptive Constraint Framework'
    tf = body_shape.text_frame
    tf.text = '''Core Hypothesis: Optimal constraint width is a function of local signal complexity

Three-Step Framework:

1. Local Complexity Quantification
   • Rolling Shannon entropy: H_i = -Σ p_k log₂(p_k)
   • High entropy → Complex regions (QRS complexes)
   • Low entropy → Simple regions (noise/isoelectric)

2. Adaptive Constraint Mapping
   • Sigmoid transformation: w_i = w_min + (w_max - w_min)/(1 + e^(-k(H_i - μ_H)))
   • w_min = 2 (rigidity), w_max = 15% (flexibility)
   • k = 2.0 (transition steepness)

3. Constrained DTW with Variable Tunnel
   • Dynamic programming with position-dependent constraints
   • Creates expanding/contracting search space
   • Prevents pathological warping while preserving morphological alignment'''

    # Add image if available
    if os.path.exists(img_paths['dtw_comparison_3']):
        left = Inches(4.0)
        top = Inches(2.0)
        width = Inches(4.0)
        height = Inches(3.0)
        slide.shapes.add_picture(img_paths['dtw_comparison_3'], left, top, width, height)

    # ===== SLIDE 5: THEORETICAL ANALYSIS =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Theoretical Foundation & Complexity'
    tf = body_shape.text_frame
    tf.text = '''Theorem: EAC-DTW strictly bounds pathological warping in low-complexity regions

Proof Sketch:
1. Low entropy regions: H_i → 0
2. Sigmoid mapping: w_i → w_min (tight constraints)
3. Geometric consequence: Path stays near diagonal
4. Noise forced to align with baseline, not morphological features

Computational Complexity:

Algorithm              Complexity    Runtime (300 samples)
Euclidean Distance     O(N)         0.4 ms
Standard DTW           O(N²)        45.2 ms
Sakoe-Chiba (10%)      O(N·R)       8.5 ms
EAC-DTW               O(N·w̄)       6.1 ms ✅

28% speedup over Sakoe-Chiba (w̄ = 8.8 < R = 36)'''

    # Add image if available
    if os.path.exists(img_paths['cost_matrix_comparison']):
        left = Inches(4.0)
        top = Inches(3.5)
        width = Inches(4.0)
        height = Inches(2.5)
        slide.shapes.add_picture(img_paths['cost_matrix_comparison'], left, top, width, height)

    # ===== SLIDE 6: EXPERIMENTAL RESULTS =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Experimental Results'
    tf = body_shape.text_frame
    tf.text = '''Dataset: Synthetic ECG-like signals (5 arrhythmia classes, 150 samples)
Evaluation: 1-NN classification with LOOCV across 3 noise conditions

Classification Accuracy:

Method              Clean    20dB SNR    10dB SNR
Euclidean Distance  92.4%    88.8%       76.5%
Standard DTW        96.1%    85.2%       68.4% ↓
Sakoe-Chiba (10%)   97.5%    91.6%       73.3%
EAC-DTW            97.8%    94.2%       79.3% ↑

Key Achievements:
✅ 6.0 percentage points improvement at 10dB SNR
✅ 41% reduction in pathological warping singularities
✅ 28% computational speedup
✅ Maintains robustness in high-noise conditions

Clinical Impact: More trustworthy automated ECG classification'''

    # Add images if available
    if os.path.exists(img_paths['EG_DTW_Performance_2025']):
        left = Inches(0.5)
        top = Inches(3.0)
        width = Inches(3.5)
        height = Inches(2.5)
        slide.shapes.add_picture(img_paths['EG_DTW_Performance_2025'], left, top, width, height)

    if os.path.exists(img_paths['singularities_comparision']):
        left = Inches(4.5)
        top = Inches(3.0)
        width = Inches(3.5)
        height = Inches(2.5)
        slide.shapes.add_picture(img_paths['singularities_comparision'], left, top, width, height)

    # ===== SLIDE 7: CONTRIBUTIONS & IMPACT =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Contributions & Business Impact'
    tf = body_shape.text_frame
    tf.text = '''Research Contributions:
✅ Novel entropy-adaptive constraint mechanism for DTW
✅ Theoretical foundation bridging rigidity-elasticity trade-off
✅ 6.0pp accuracy improvement in noisy ECG conditions
✅ 41% reduction in pathological warping instances

Business Impact:
• Healthcare Market: $6.2B ECG market (growing to $9.8B by 2030)
• Cost Savings: $2,500-5,000 per false alarm avoided
• Applications: Hospital ECG systems, wearable monitors, telemedicine
• Competitive Advantage: Noise-robust, computationally efficient, interpretable

Limitations:
⚠️ Synthetic data evaluation (clinical validation needed)
⚠️ Single-lead analysis (multi-lead extension planned)

Future Work: MIT-BIH clinical validation, real-time implementation, multi-lead extension'''

    # ===== SLIDE 8: CONCLUSION & Q&A =====
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Conclusion & Q&A'
    tf = body_shape.text_frame
    tf.text = '''Summary:
EAC-DTW addresses the fundamental limitation of fixed-constraint DTW by using local signal entropy to dynamically adjust warping flexibility. Our approach achieves superior noise robustness while maintaining computational efficiency, making it suitable for clinical ECG analysis.

Key Results:
• 79.3% accuracy at 10dB SNR (+6.0pp improvement)
• 41% reduction in pathological warping
• 28% computational speedup
• Foundation for trustworthy automated cardiac monitoring

Questions & Discussion:

We welcome questions about:
• Technical implementation details
• Algorithm performance characteristics
• Clinical validation plans
• Business development opportunities

Contact:
Fnu Ashutosh (an05893n@pace.edu)
Shivam Jha (sj34101n@pace.edu)'''

    # Save the presentation
    output_path = os.path.join(os.getcwd(), 'EAC_DTW_Condensed_Presentation.pptx')
    prs.save(output_path)
    print(f"Condensed presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_condensed_presentation()